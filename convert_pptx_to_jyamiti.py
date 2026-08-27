#!/usr/bin/env python3
"""
High-Fidelity PPTX to Jyamiti Custom SlideDeck Converter
=========================================================
Converts PowerPoint (.pptx) presentation files into Jyamiti-compatible
high-fidelity SlideDeck JSON files (Option 1 Hybrid Approach).

Features:
1. 100% Identical Visual Rendering:
   - Uses Windows PowerPoint COM or LibreOffice to export 1080p high-res slide images.
2. Content & Notes Extraction:
   - Extracts slide titles, text frames, bullet points, and speaker notes.
3. Jyamiti Schema Compliance:
   - Outputs ready-to-import SlideDeck JSON compatible with Jyamiti MongoDB & Flutter CMS.
4. Embedded Base64 or Image Asset Paths:
   - Optionally embeds images directly as data URIs (for instant standalone import)
     or exports to a companion assets directory.

Usage:
  python convert_pptx_to_jyamiti.py path/to/presentation.pptx
  python convert_pptx_to_jyamiti.py path/to/presentation.pptx --course "Physics 101" --embed-base64
  python convert_pptx_to_jyamiti.py ./my_pptx_folder/ --batch
"""

import os
import sys
import json
import argparse
import subprocess
from pathlib import Path
import base64
import time

def export_slides_via_powerpoint_com(pptx_path: Path, output_img_dir: Path, width: int = 1920, height: int = 1080):
    """
    Exports slides as PNG images using Windows PowerPoint COM (via PowerShell).
    Works on any Windows system with Microsoft PowerPoint installed.
    """
    output_img_dir.mkdir(parents=True, exist_ok=True)
    abs_pptx = str(pptx_path.resolve()).replace("'", "''")
    abs_out = str(output_img_dir.resolve()).replace("'", "''")

    ps_script = f"""
    $ppApp = New-Object -ComObject PowerPoint.Application
    $ppApp.Visible = [Microsoft.Office.Core.MsoTriState]::msoFalse
    try {{
        $pres = $ppApp.Presentations.Open('{abs_pptx}', [Microsoft.Office.Core.MsoTriState]::msoTrue, [Microsoft.Office.Core.MsoTriState]::msoFalse, [Microsoft.Office.Core.MsoTriState]::msoFalse)
        $slideIndex = 0
        foreach ($slide in $pres.Slides) {{
            $slideIndex++
            $imgName = "slide_$slideIndex.png"
            $imgPath = Join-Path '{abs_out}' $imgName
            $slide.Export($imgPath, 'PNG', {width}, {height})
        }}
        $pres.Close()
        Write-Output "SUCCESS:$slideIndex"
    }} finally {{
        $ppApp.Quit()
        [System.GC]::Collect()
        [System.GC]::WaitForPendingFinalizers()
    }}
    """

    try:
        res = subprocess.run(
            ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", ps_script],
            capture_output=True,
            text=True,
            timeout=180
        )
        if "SUCCESS:" in res.stdout:
            lines = [l for l in res.stdout.splitlines() if "SUCCESS:" in l]
            count = int(lines[0].split("SUCCESS:")[1].strip())
            return count
        else:
            print(f"[COM Warning] PowerPoint COM returned: {res.stderr.strip() or res.stdout.strip()}")
            return 0
    except Exception as e:
        print(f"[COM Error] PowerPoint COM export failed: {e}")
        return 0

def export_slides_via_libreoffice(pptx_path: Path, output_img_dir: Path):
    """
    Fallback converter using LibreOffice + pdftoppm / PyMuPDF.
    """
    output_img_dir.mkdir(parents=True, exist_ok=True)
    pdf_dir = output_img_dir / "temp_pdf"
    pdf_dir.mkdir(exist_ok=True)
    
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", str(pptx_path), "--outdir", str(pdf_dir)],
            check=True,
            capture_output=True,
            timeout=120
        )
        pdf_file = next(pdf_dir.glob("*.pdf"), None)
        if not pdf_file:
            return 0
        
        # Convert PDF pages to PNG
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(pdf_file)
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                pix.save(str(output_img_dir / f"slide_{i+1}.png"))
            doc.close()
            return len(doc)
        except ImportError:
            # Try pdftoppm
            subprocess.run(
                ["pdftoppm", "-png", "-r", "150", str(pdf_file), str(output_img_dir / "slide")],
                check=True
            )
            return len(list(output_img_dir.glob("slide*.png")))
    except Exception as e:
        print(f"[LibreOffice Warning] LibreOffice conversion not available: {e}")
        return 0
    finally:
        import shutil
        if pdf_dir.exists():
            shutil.rmtree(pdf_dir, ignore_errors=True)

def extract_pptx_metadata_and_text(pptx_path: Path):
    """
    Extracts text and titles from PPTX using python-pptx or ZIP XML parser.
    """
    slides_info = []
    
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        for idx, slide in enumerate(prs.slides):
            title = ""
            texts = []
            notes = ""
            
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text.strip()
                    if not t:
                        continue
                    if shape == slide.shapes.title or (not title and len(t) < 80 and '\n' not in t):
                        title = t
                    else:
                        texts.append(t)
                        
            if not title:
                title = f"Slide {idx + 1}"
                
            slides_info.append({
                "slideIndex": idx,
                "title": title,
                "paragraphs": texts,
                "notes": notes,
            })
        return slides_info
    except ImportError:
        pass

    # Fallback: simple ZIP inspection
    import zipfile
    import xml.etree.ElementTree as ET
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            slide_entries = [f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
            slide_entries.sort(key=lambda x: int(''.join(filter(str.isdigit, x)) or 0))
            
            for idx, entry in enumerate(slide_entries):
                xml_content = z.read(entry)
                root = ET.fromstring(xml_content)
                texts = []
                for elem in root.iter():
                    if elem.tag.endswith('}t') and elem.text:
                        texts.append(elem.text.strip())
                
                combined = " ".join(texts)
                title = texts[0] if texts else f"Slide {idx + 1}"
                slides_info.append({
                    "slideIndex": idx,
                    "title": title[:60],
                    "paragraphs": [combined] if combined else [],
                    "notes": "",
                })
        return slides_info
    except Exception as e:
        print(f"[Extraction Warning] Text extraction error: {e}")
        return []

def convert_single_pptx(
    pptx_file: Path,
    output_dir: Path,
    course_id: str = "course_101",
    course_name: str = "Imported Course",
    embed_base64: bool = False,
    theme: str = "darkGlass",
):
    print(f"\n=======================================================")
    print(f" Converting: {pptx_file.name}")
    print(f"=======================================================")

    output_dir.mkdir(parents=True, exist_ok=True)
    images_dir = output_dir / f"{pptx_file.stem}_slides"
    
    # 1. Render slide images
    print("-> Rendering 1:1 High-Fidelity Slide Images...")
    slide_count = export_slides_via_powerpoint_com(pptx_file, images_dir)
    if slide_count == 0:
        slide_count = export_slides_via_libreoffice(pptx_file, images_dir)
        
    # 2. Extract text metadata
    print("-> Extracting Slide Titles and Content...")
    meta_slides = extract_pptx_metadata_and_text(pptx_file)
    
    total_slides = max(slide_count, len(meta_slides), 1)
    print(f"-> Total Slides Processed: {total_slides}")
    
    deck_id = f"deck_{int(time.time() * 1000)}"
    deck_title = pptx_file.stem.replace("_", " ").replace("-", " ").title()
    
    slide_items = []
    
    for i in range(total_slides):
        meta = meta_slides[i] if i < len(meta_slides) else {}
        slide_title = meta.get("title") or f"Slide {i + 1}"
        
        # Locate corresponding slide image
        img_path = images_dir / f"slide_{i + 1}.png"
        img_url = ""
        
        if img_path.exists():
            if embed_base64:
                with open(img_path, "rb") as img_f:
                    b64 = base64.b64encode(img_f.read()).decode("utf-8")
                    img_url = f"data:image/png;base64,{b64}"
            else:
                # Relative local path / URL
                img_url = str(img_path.resolve())

        # Construct optional blocks from extracted text
        blocks = []
        paragraphs = meta.get("paragraphs", [])
        if paragraphs:
            for p_idx, para in enumerate(paragraphs[:3]):  # Keep first few key points
                blocks.append({
                    "id": f"b_{i}_{p_idx}",
                    "type": "paragraph",
                    "content": para,
                    "extra": None,
                    "caption": None,
                })
        
        slide_items.append({
            "id": f"s_{deck_id}_{i}",
            "slideIndex": i,
            "title": slide_title,
            "theme": theme,
            "imageUrl": img_url,
            "enableWhiteboard": True,
            "blocks": blocks,
            "quiz": None,
        })
        
    deck_data = {
        "id": deck_id,
        "courseId": course_id,
        "courseName": course_name,
        "title": deck_title,
        "description": f"Imported high-fidelity presentation converted from {pptx_file.name}.",
        "createdAt": time.strftime("%Y-%m-%dT%H:%M:%S.000Z", time.gmtime()),
        "isPublished": True,
        "isDownloadedOffline": True,
        "slides": slide_items,
    }
    
    out_json = output_dir / f"{pptx_file.stem}_deck.json"
    with open(out_json, "w", encoding="utf-8") as f:
        json.dump(deck_data, f, indent=2, ensure_ascii=False)
        
    print(f"\n [SUCCESS] Converted deck saved to: {out_json}")
    print(f" -> Slides: {len(slide_items)}")
    print(f" -> You can directly import this JSON file into the Jyamiti App via Slide Decks CMS!")
    return out_json

def main():
    parser = argparse.ArgumentParser(description="Convert PPTX to Jyamiti Custom High-Fidelity SlideDecks")
    parser.add_argument("input_path", help="Path to .pptx file or folder containing .pptx files")
    parser.add_argument("--output", "-o", default="./converted_slides", help="Output directory (default: ./converted_slides)")
    parser.add_argument("--course", "-c", default="Mathematics", help="Course Name")
    parser.add_argument("--course-id", default="course_101", help="Course ID")
    parser.add_argument("--theme", default="darkGlass", help="Slide Theme (darkGlass, midnightNeon, cleanLight, etc.)")
    parser.add_argument("--embed-base64", action="store_true", help="Embed slide images as Base64 strings into the JSON")
    
    args = parser.parse_args()
    input_p = Path(args.input_path)
    output_p = Path(args.output)
    
    if not input_p.exists():
        print(f"[Error] Path does not exist: {input_p}")
        sys.exit(1)
        
    if input_p.is_dir():
        pptx_files = list(input_p.glob("*.pptx")) + list(input_p.glob("*.ppt"))
        if not pptx_files:
            print(f"No .pptx files found in directory {input_p}")
            sys.exit(0)
        print(f"Found {len(pptx_files)} presentation(s) to convert.")
        for f in pptx_files:
            convert_single_pptx(f, output_p, args.course_id, args.course, args.embed_base64, args.theme)
    else:
        convert_single_pptx(input_p, output_p, args.course_id, args.course, args.embed_base64, args.theme)

if __name__ == "__main__":
    main()
